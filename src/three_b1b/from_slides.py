"""from-slides: extract PPTX content, plan video, generate Manim animation."""

from __future__ import annotations

import base64
import sys
from pathlib import Path
from typing import NoReturn, Optional

import click

from .generate import (
    PROVIDERS,
    _build_system_prompt,
    _extract_code,
    _render_scene,
    _resolve_api_key,
    call_llm,
)

VISION_PROVIDERS: set[str] = {"anthropic", "openai"}
EXTRACTION_MODELS: dict[str, str] = {
    "anthropic": "claude-haiku-4-5-20251001",
    "openai": "gpt-4o-mini",
}
PLAN_MODELS: dict[str, str] = {
    "anthropic": "claude-haiku-4-5-20251001",
    "openai": "gpt-4o-mini",
    "claude-code": "opus",
}

_AVG_TEXT_THRESHOLD = 80
_PIC_RATIO_THRESHOLD = 0.35

from .prompts import IMAGE_DESCRIBE as _IMAGE_PROMPT, SLIDES_PLAN, SLIDES_GENERATE


def _pptx_missing() -> NoReturn:
    click.echo("python-pptx is required: pip install '3brown1blue[slides]'")
    sys.exit(1)


def _load_deck(deck: Path) -> tuple[list[dict], float, float]:
    """Open the PPTX once and return slides data, avg text per slide, and pic ratio."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        _pptx_missing()

    prs = Presentation(deck)
    slides_data: list[dict] = []
    total_shapes = pic_count = total_chars = 0

    for idx, slide in enumerate(prs.slides, 1):
        title = ""
        texts: list[str] = []
        images: list[tuple[bytes, str]] = []

        for shape in slide.shapes:
            total_shapes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic_count += 1
                try:
                    images.append((shape.image.blob, shape.image.content_type))
                except Exception:
                    pass  # skip linked (non-embedded) images
            elif shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if not t:
                    continue
                if not title and shape.name.lower().startswith("title"):
                    title = t
                texts.append(t)
                total_chars += len(t)

        if not title and texts:
            title = texts[0]

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append({
            "index": idx, "title": title,
            "text": "\n".join(texts), "notes": notes, "images": images,
        })

    n = len(slides_data)
    avg_text = total_chars / n if n else 0.0
    pic_ratio = pic_count / total_shapes if total_shapes else 0.0
    return slides_data, avg_text, pic_ratio


def _detect_mode(avg_text: float, pic_ratio: float, provider: str) -> str:
    needs_vision = avg_text < _AVG_TEXT_THRESHOLD or pic_ratio > _PIC_RATIO_THRESHOLD
    if needs_vision and provider not in VISION_PROVIDERS:
        click.echo(
            f"Vision extraction needs anthropic or openai; {provider} detected"
            " -- falling back to text extraction."
        )
        return "text"
    return "vision" if needs_vision else "text"


def _describe_images(
    images: list[tuple[bytes, str]], provider: str, model: str, api_key: str
) -> list[str]:
    """Describe embedded slide images using a cheap multimodal model."""
    if provider not in VISION_PROVIDERS:
        raise ValueError(f"_describe_images called with non-vision provider: {provider!r}")
    if provider == "anthropic":
        try:
            import anthropic
        except ImportError:
            click.echo("Run: pip install '3brown1blue[anthropic]'")
            sys.exit(1)
        client = anthropic.Anthropic(api_key=api_key, timeout=60.0)
    else:
        try:
            from openai import OpenAI
        except ImportError:
            click.echo("Run: pip install '3brown1blue[openai]'")
            sys.exit(1)
        client = OpenAI(api_key=api_key, timeout=60.0)  # type: ignore[assignment]

    descriptions: list[str] = []
    for blob, mime in images:
        b64 = base64.b64encode(blob).decode()
        try:
            if provider == "anthropic":
                msg = client.messages.create(  # type: ignore[union-attr]
                    model=model, max_tokens=512,
                    messages=[{"role": "user", "content": [
                        {"type": "image", "source": {"type": "base64", "media_type": mime, "data": b64}},
                        {"type": "text", "text": _IMAGE_PROMPT},
                    ]}],
                )
                descriptions.append(msg.content[0].text)
            else:
                resp = client.chat.completions.create(  # type: ignore[union-attr]
                    model=model, max_tokens=512,
                    messages=[{"role": "user", "content": [
                        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                        {"type": "text", "text": _IMAGE_PROMPT},
                    ]}],
                )
                descriptions.append(resp.choices[0].message.content or "")
        except Exception as e:
            click.echo(f"  Image description failed: {e}")
            descriptions.append("[unavailable]")
    return descriptions


def _format_slides_markdown(slides: list[dict], vision_descs: dict[int, list[str]]) -> str:
    parts: list[str] = []
    for s in slides:
        header = f"## Slide {s['index']}: {s['title']}" if s["title"] else f"## Slide {s['index']}"
        body: list[str] = []
        if s["text"]:
            body.append(s["text"])
        for i, desc in enumerate(vision_descs.get(s["index"], []), 1):
            body.append(f"[Image {i}]: {desc}")
        if s["notes"]:
            body.append(f"[Speaker notes]: {s['notes']}")
        parts.append(header + "\n" + "\n".join(body))
    return "\n\n".join(parts)


@click.command(name="from-slides")
@click.argument("deck", type=click.Path(exists=True, dir_okay=False, path_type=Path))
@click.option("--provider", "-p", type=click.Choice(list(PROVIDERS)), default=None,
              help="LLM provider for planning and generation.")
@click.option("--model", "-m", default=None, help="Generation model (provider default if omitted).")
@click.option("--plan-model", default=None, help="Planning model (haiku/gpt-4o-mini by default).")
@click.option("--api-key", "-k", default=None, help="API key (or set the provider env var).")
@click.option("--output", "-o", default="scene.py", show_default=True, help="Output file.")
@click.option("--plan-output", default=None, help="Optional markdown file to save the reviewed plan.")
@click.option("--render", is_flag=True, help="Run manim after generation.")
@click.option("--quality", "-q", type=click.Choice(["l", "m", "h", "k"]), default="l",
              show_default=True, help="Render quality: l=fast, m=medium, h=1080p, k=4K.")
@click.option("--mode", type=click.Choice(["auto", "text", "vision"]), default="auto",
              show_default=True, help="Slide extraction mode (auto detects from content).")
@click.option("--audience", "-a",
              type=click.Choice(["high-school", "undergrad", "graduate", "industry"]),
              default="undergrad", show_default=True,
              help="Target audience level.")
@click.option("--domain", "-d",
              type=click.Choice(["auto", "machine-learning", "mathematics", "physics", "biology", "security", "neuroscience", "algorithms"]),
              default="auto", show_default=True,
              help="Academic domain (auto-detect if omitted).")
def from_slides(
    deck: Path,
    provider: Optional[str],
    model: Optional[str],
    plan_model: Optional[str],
    api_key: Optional[str],
    output: str,
    plan_output: Optional[str],
    render: bool,
    quality: str,
    mode: str,
    audience: str,
    domain: str,
) -> None:
    """Generate a Manim animated explainer from a PowerPoint file.

    Extracts slide content, plans the video with a fast model, lets you
    review and edit the plan, then generates complete Manim code.

    \b
    Examples:
      3brown1blue from-slides lecture.pptx --provider anthropic --render
      3brown1blue from-slides talk.pptx -p openai --mode vision -o talk.py
      3brown1blue from-slides lecture.pptx -p anthropic -a graduate -d neuroscience
    """
    from ._shared import prompt_provider, confirm_plan

    if provider is None:
        provider = prompt_provider()

    cfg = PROVIDERS[provider]
    model = model or cfg["default_model"]
    plan_model = plan_model or PLAN_MODELS.get(provider, cfg["default_model"])
    api_key = _resolve_api_key(provider, api_key)

    click.echo(f"Analyzing {deck.name}...")
    slides, avg_text, pic_ratio = _load_deck(deck)
    click.echo(f"  {len(slides)} slides detected.")

    effective_mode = mode
    if effective_mode == "auto":
        effective_mode = _detect_mode(avg_text, pic_ratio, provider)
    elif effective_mode == "vision" and provider not in VISION_PROVIDERS:
        click.echo("Vision mode requires anthropic or openai -- falling back to text.")
        effective_mode = "text"
    click.echo(f"  Extraction mode: {effective_mode}")

    vision_descs: dict[int, list[str]] = {}
    if effective_mode == "vision":
        ext_model = EXTRACTION_MODELS.get(provider, PROVIDERS[provider]["default_model"])
        total_images = sum(len(s["images"]) for s in slides)
        if total_images:
            click.echo(f"  Describing {total_images} images with {ext_model}...")
            for slide in slides:
                if slide["images"]:
                    vision_descs[slide["index"]] = _describe_images(
                        slide["images"], provider, ext_model, api_key
                    )

    slides_md = _format_slides_markdown(slides, vision_descs)
    system = _build_system_prompt(audience=audience, domain=domain)

    click.echo(f"\nPlanning video with {plan_model}...")
    click.echo(f"  Audience: {audience}  Domain: {domain}")
    plan = call_llm(provider, plan_model, api_key,
                    SLIDES_PLAN.format(slides_content=slides_md, audience=audience, domain=domain),
                    system=system, audience=audience, domain=domain)
    plan = confirm_plan(plan)
    if plan_output:
        plan_path = Path(plan_output)
        plan_path.write_text(plan)
        click.echo(f"Saved plan: {plan_path}")

    click.echo(f"\nGenerating Manim code with {model}...")
    raw = call_llm(provider, model, api_key, SLIDES_GENERATE.format(plan=plan),
                   system=system, audience=audience, domain=domain)
    code = _extract_code(raw)

    out_path = Path(output)
    out_path.write_text(code)
    click.echo(f"Saved:  {out_path}")

    if render:
        _render_scene(out_path, quality)
